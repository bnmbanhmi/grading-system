#!/usr/bin/env python3
"""
Simple script to run the assignment organizer on the course zip file.
Edit the ZIP_FILE_PATH variable below to point to your zip file.
"""

import os
import zipfile
import shutil
import glob
import re
import subprocess
import sys
from pathlib import Path
from typing import List, Optional, Dict


class AssignmentOrganizer:
    def __init__(self, base_path: str):
        self.base_path = Path(base_path)
        self.extracted_folders = []
        
    def extract_group_zips(self, group_folder: Path) -> None:
        """Recursively extract all zip files within a group folder."""
        print(f"Processing group folder: {group_folder.name}")
        
        # Recursively extract all zip files until no more are found
        extracted_count = 0
        max_iterations = 10  # Prevent infinite loops
        iteration = 0
        
        while iteration < max_iterations:
            # Find all zip files in the group folder (including subdirectories)
            zip_files = list(group_folder.rglob("*.zip"))
            
            if not zip_files:
                break  # No more zip files to extract
                
            iteration += 1
            print(f"  Extraction iteration {iteration}: Found {len(zip_files)} zip files")
            
            for zip_file in zip_files:
                print(f"    Extracting: {zip_file.name}")
                try:
                    # Create a temporary extraction directory
                    temp_extract_dir = zip_file.parent / f"temp_extract_{zip_file.stem}"
                    temp_extract_dir.mkdir(exist_ok=True)
                    
                    with zipfile.ZipFile(zip_file, 'r') as zip_ref:
                        # Extract to temporary directory
                        zip_ref.extractall(temp_extract_dir)
                    
                    # Move all contents from temp directory to the zip file's location
                    self.move_extracted_contents(temp_extract_dir, zip_file.parent)
                    
                    # Clean up temp directory
                    shutil.rmtree(temp_extract_dir)
                    
                    # Delete the zip file after extraction
                    zip_file.unlink()
                    print(f"    Deleted: {zip_file.name}")
                    extracted_count += 1
                    
                except zipfile.BadZipFile:
                    print(f"    Warning: {zip_file.name} is not a valid zip file")
                except Exception as e:
                    print(f"    Error extracting {zip_file.name}: {e}")
        
        if extracted_count > 0:
            print(f"  Total extracted: {extracted_count} zip files in {iteration} iterations")
        else:
            print(f"  No zip files found in {group_folder.name}")
    
    def move_extracted_contents(self, source_dir: Path, target_dir: Path) -> None:
        """Move all contents from source directory to target directory."""
        for item in source_dir.rglob("*"):
            if item.is_file():
                # Calculate relative path from source_dir
                relative_path = item.relative_to(source_dir)
                target_path = target_dir / relative_path
                
                # Create parent directories if they don't exist
                target_path.parent.mkdir(parents=True, exist_ok=True)
                
                # Handle file name conflicts
                counter = 1
                original_target = target_path
                while target_path.exists():
                    name_part = original_target.stem
                    suffix_part = original_target.suffix
                    target_path = original_target.parent / f"{name_part}_{counter}{suffix_part}"
                    counter += 1
                
                # Move the file
                shutil.move(str(item), str(target_path))
                print(f"    Moved: {relative_path} → {target_path.name}")
    
    def find_source_code_items(self, directory: Path) -> List[Path]:
        """Find source code files and folders."""
        source_items = []
        
        # Common patterns for source code
        code_patterns = [
            "*code*", "*script*", "*src*", "*source*",
            "*.cs", "*.py", "*.js", "*.cpp", "*.java", "*.c",
            "*.h", "*.hpp", "*.unity", "*.meta"
        ]
        
        for pattern in code_patterns:
            items = list(directory.glob(pattern))
            items.extend(directory.rglob(pattern))
            source_items.extend(items)
        
        # Remove duplicates and ensure they're within the directory
        unique_items = []
        for item in source_items:
            if item not in unique_items and directory in item.parents or item.parent == directory:
                unique_items.append(item)
        
        return unique_items
    
    def identify_file_types(self, directory: Path) -> Dict[str, List[Path]]:
        """Categorize files by type."""
        file_types = {
            'pdf': [],
            'excel': [],
            'video': [],
            'presentation': [],
            'source_code': [],
            'other': []
        }
        
        # Define file extensions for each category
        extensions = {
            'pdf': ['.pdf'],
            'excel': ['.xlsx', '.xls', '.csv'],
            'video': ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv'],
            'presentation': ['.ppt', '.pptx', '.key']
        }
        
        # Get all files in directory and subdirectories
        all_files = []
        for item in directory.rglob("*"):
            if item.is_file() and not item.name.startswith('.'):
                all_files.append(item)
        
        # Categorize files
        for file_path in all_files:
            categorized = False
            
            for category, exts in extensions.items():
                if file_path.suffix.lower() in exts:
                    file_types[category].append(file_path)
                    categorized = True
                    break
            
            if not categorized:
                # Check if it's source code
                if self.is_source_code_file(file_path):
                    file_types['source_code'].append(file_path)
                else:
                    file_types['other'].append(file_path)
        
        return file_types
    
    def is_source_code_file(self, file_path: Path) -> bool:
        """Check if a file is likely source code that's useful for assessment."""
        
        # Only include actual source code files (exclude Unity metadata and binary files)
        assessment_useful_extensions = {
            '.cs',      # C# scripts (Unity/general)
            '.py',      # Python scripts
            '.js',      # JavaScript
            '.ts',      # TypeScript
            '.cpp',     # C++
            '.c',       # C
            '.h',       # Header files
            '.hpp',     # C++ headers
            '.java',    # Java
            '.kt',      # Kotlin
            '.swift',   # Swift
            '.go',      # Go
            '.rs',      # Rust
            '.php',     # PHP
            '.rb',      # Ruby
            '.shader',  # Shader files (actual code)
            '.hlsl',    # HLSL shaders
            '.glsl',    # GLSL shaders
            '.json',    # Configuration/data files
            '.xml',     # Configuration files
            '.yaml',    # Configuration files
            '.yml'      # Configuration files
        }
        
        # Exclude Unity metadata and binary files that don't contribute to assessment
        excluded_extensions = {
            '.meta',    # Unity metadata
            '.prefab',  # Unity prefabs (binary)
            '.unity',   # Unity scene files (binary) 
            '.mat',     # Unity materials (binary)
            '.asset',   # Unity asset files (binary)
            '.anim',    # Unity animations (binary)
            '.controller', # Unity animator controllers (binary)
        }
        
        # Exclude common non-code files
        excluded_filenames = {
            '.ds_store', 'thumbs.db', 'desktop.ini'
        }
        
        file_ext = file_path.suffix.lower()
        filename_lower = file_path.name.lower()
        
        # Exclude explicitly excluded extensions
        if file_ext in excluded_extensions:
            return False
            
        # Exclude common system files
        if filename_lower in excluded_filenames:
            return False
        
        # Include files with useful extensions
        if file_ext in assessment_useful_extensions:
            return True
        
        # Check if filename contains code-related keywords for extensionless files
        code_keywords = ['script', 'code', 'src', 'source']
        filename_lower = file_path.name.lower()
        return any(keyword in filename_lower for keyword in code_keywords)
    
    def rename_with_group_prefix(self, file_path: Path, group_name: str) -> str:
        """Generate new filename with group prefix."""
        # Remove common prefixes that might already exist
        name = file_path.stem
        suffix = file_path.suffix
        
        # Clean up existing group prefixes
        name = re.sub(r'^(GC\d+|OL\d+)-?', '', name, flags=re.IGNORECASE)
        
        # Add group prefix
        new_name = f"{group_name}-{name}{suffix}"
        return new_name
    
    def organize_group_folder(self, group_folder: Path, organized_base_path: Path) -> None:
        """Organize a single group's assignment folder."""
        group_name = group_folder.name
        print(f"\nOrganizing group: {group_name}")
        
        # Create "organized_assignments" directory if it doesn't exist
        organized_assignments_dir = organized_base_path / "organized_assignments"
        organized_assignments_dir.mkdir(exist_ok=True)
        
        # Create separate organized folder for this group in the organized_assignments directory
        organized_path = organized_assignments_dir / f"{group_name}_organized"
        
        # Remove existing organized folder if it exists (for clean re-organization)
        if organized_path.exists():
            shutil.rmtree(organized_path)
        
        organized_path.mkdir(exist_ok=True)
        
        # Create subdirectories (removed presentations folder)
        subdirs = {
            'documents': organized_path / "documents",
            'videos': organized_path / "videos", 
            'source_code': organized_path / "source_code"
        }
        
        for subdir in subdirs.values():
            subdir.mkdir(exist_ok=True)
        
        # Detect project type and generate report
        project_type = self.detect_project_type(group_folder)
        print(f"  Detected project type: {project_type}")
        
        # Identify and categorize files
        file_types = self.identify_file_types(group_folder)
        
        # Move and rename files
        self.move_categorized_files(file_types, subdirs, group_name, group_folder)
        
        # Handle source code folders
        self.handle_source_code_folders(group_folder, subdirs['source_code'], group_name)
        
        # Generate organization report
        self.generate_organization_report(organized_path, group_name, file_types, project_type)
        
        print(f"✓ Organized {group_name} → organized_assignments/{organized_path.name}")
    
    def move_categorized_files(self, file_types: Dict[str, List[Path]], 
                             subdirs: Dict[str, Path], group_name: str, 
                             source_folder: Path) -> None:
        """Move and rename categorized files, converting presentations to PDF."""
        
        # Mapping of file types to destination folders (presentations now go to documents)
        destination_mapping = {
            'pdf': subdirs['documents'],
            'excel': subdirs['documents'],
            'video': subdirs['videos']
        }
        
        # Handle regular files (PDFs, Excel, Videos)
        for file_type, files in file_types.items():
            if file_type in destination_mapping and files:
                dest_folder = destination_mapping[file_type]
                
                for file_path in files:
                    # Skip files that are already in the organized folder
                    if 'organized' in str(file_path):
                        continue
                        
                    new_name = self.rename_with_group_prefix(file_path, group_name)
                    dest_path = dest_folder / new_name
                    
                    try:
                        # Ensure we don't overwrite existing files
                        counter = 1
                        while dest_path.exists():
                            name_part = dest_path.stem
                            suffix_part = dest_path.suffix
                            dest_path = dest_folder / f"{name_part}_{counter}{suffix_part}"
                            counter += 1
                        
                        shutil.copy2(file_path, dest_path)
                        print(f"  Copied: {file_path.name} → {dest_path.name}")
                        
                    except Exception as e:
                        print(f"  Error copying {file_path.name}: {e}")
        
        # Handle presentations - convert to PDF and put in documents folder
        if 'presentation' in file_types and file_types['presentation']:
            documents_folder = subdirs['documents']
            
            for presentation_path in file_types['presentation']:
                # Skip files that are already in the organized folder
                if 'organized' in str(presentation_path):
                    continue
                
                print(f"  Converting presentation: {presentation_path.name}")
                try:
                    # Convert presentation to PDF and save in documents folder
                    converted_file = self.convert_presentation_to_pdf(
                        presentation_path, documents_folder, group_name
                    )
                    
                except Exception as e:
                    print(f"  Error processing presentation {presentation_path.name}: {e}")
                    # Fallback: copy original file to documents folder with group prefix
                    try:
                        fallback_name = self.rename_with_group_prefix(presentation_path, group_name)
                        fallback_path = documents_folder / fallback_name
                        
                        counter = 1
                        while fallback_path.exists():
                            name_part = fallback_path.stem
                            suffix_part = fallback_path.suffix
                            fallback_path = documents_folder / f"{name_part}_{counter}{suffix_part}"
                            counter += 1
                        
                        shutil.copy2(presentation_path, fallback_path)
                        print(f"  Copied original: {presentation_path.name} → {fallback_path.name}")
                    except Exception as copy_error:
                        print(f"  Failed to copy original presentation: {copy_error}")
    
    def handle_source_code_folders(self, group_folder: Path, 
                                 source_dest: Path, group_name: str) -> None:
        """Handle source code folders and files, copying only assessment-useful files."""
        
        # Look for folders that might contain source code
        code_folder_patterns = ['*script*', '*code*', '*src*', '*source*']
        code_folders = []
        
        for pattern in code_folder_patterns:
            folders = list(group_folder.glob(pattern))
            code_folders.extend([f for f in folders if f.is_dir() and 'organized' not in str(f)])
        
        # Copy source code folders selectively (excluding .meta files and other unnecessary files)
        for folder in code_folders:
            dest_folder_name = f"{group_name}-{folder.name}"
            dest_path = source_dest / dest_folder_name
            
            try:
                if dest_path.exists():
                    shutil.rmtree(dest_path)
                dest_path.mkdir(parents=True, exist_ok=True)
                
                # Copy only useful files
                self.copy_source_code_selectively(folder, dest_path, group_name)
                print(f"  Copied source folder: {folder.name} → {dest_folder_name}")
            except Exception as e:
                print(f"  Error copying source folder {folder.name}: {e}")
        
        # Copy individual assessment-useful source files
        individual_files = []
        for item in group_folder.rglob("*"):
            if (item.is_file() and 
                self.is_source_code_file(item) and 
                'organized' not in str(item) and
                not any(str(item).startswith(str(folder)) for folder in code_folders)):
                individual_files.append(item)
        
        if individual_files:
            misc_folder = source_dest / f"{group_name}-misc-source"
            misc_folder.mkdir(exist_ok=True)
            
            for file_path in individual_files:
                try:
                    dest_file = misc_folder / file_path.name
                    # Handle name conflicts
                    counter = 1
                    while dest_file.exists():
                        stem = file_path.stem
                        suffix = file_path.suffix
                        dest_file = misc_folder / f"{stem}_{counter}{suffix}"
                        counter += 1
                    
                    shutil.copy2(file_path, dest_file)
                except Exception as e:
                    print(f"  Error copying individual file {file_path.name}: {e}")
            
            print(f"  Collected {len(individual_files)} individual source files → {group_name}-misc-source")
        
        # Clean up any existing unnecessary files in the destination
        self.clean_unnecessary_files(source_dest)
    
    def process_course_zip(self, zip_file_path: str) -> None:
        """Main method to process a course assignment zip file."""
        zip_path = Path(zip_file_path)
        
        if not zip_path.exists():
            print(f"Error: Zip file {zip_file_path} not found")
            return
        
        print(f"Processing course zip: {zip_path.name}")
        
        # Extract the main zip file
        extract_dir = zip_path.parent / zip_path.stem
        extract_dir.mkdir(exist_ok=True)
        
        try:
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
            print(f"Extracted main zip to: {extract_dir}")
        except Exception as e:
            print(f"Error extracting main zip: {e}")
            return
        
        # Find the directory that actually contains group folders (handles nested structures)
        course_dir = find_group_folders_directory(str(extract_dir))
        course_path = Path(course_dir)
        
        # Find group folders using the robust pattern matching
        all_subdirs = [d for d in course_path.iterdir() 
                      if d.is_dir() and not d.name.startswith('.')]
        
        group_folders = [d for d in all_subdirs if is_group_folder(d.name)]
        non_group_folders = [d for d in all_subdirs if not is_group_folder(d.name)]
        
        print(f"\nUsing course directory: {course_dir}")
        print(f"Found {len(group_folders)} group folders:")
        for folder in group_folders:
            print(f"  - {folder.name}")
        
        if non_group_folders:
            print(f"Found {len(non_group_folders)} non-group folders (ignored):")
            for folder in non_group_folders:
                print(f"  - {folder.name}")
        
        # Process each group folder
        for group_folder in group_folders:
            self.extract_group_zips(group_folder)
        
        print(f"\n✅ Extraction complete! Processed {len(group_folders)} groups.")
        print(f"Files extracted to: {course_dir}")
        
        # Now organize the extracted files
        self.organize_all_groups(course_path)
        print("Each group folder now contains all extracted files from any zip files that were inside.")
    
    def convert_presentation_to_pdf(self, presentation_path: Path, output_dir: Path, group_name: str) -> Path:
        """
        Convert presentation files to PDF format.
        
        Args:
            presentation_path (Path): Path to the presentation file
            output_dir (Path): Directory to save the converted PDF
            group_name (str): Group name for file naming
            
        Returns:
            Path: Path to the converted PDF file, or original file path if conversion fails
        """
        try:
            # Generate output PDF filename
            pdf_filename = f"{group_name}-{presentation_path.stem}.pdf"
            pdf_output_path = output_dir / pdf_filename
            
            # Ensure output directory exists
            output_dir.mkdir(parents=True, exist_ok=True)
            
            file_ext = presentation_path.suffix.lower()
            
            # Try different conversion methods based on the system
            success = False
            
            # Method 1: Try LibreOffice (cross-platform)
            if not success:
                success = self._convert_with_libreoffice(presentation_path, output_dir, pdf_output_path)
            
            # Method 2: Try using Python-pptx for basic conversion (fallback)
            if not success and file_ext in ['.pptx', '.ppt']:
                success = self._convert_with_python_libs(presentation_path, pdf_output_path)
            
            # Method 3: macOS-specific conversion using Keynote/PowerPoint
            if not success and sys.platform == 'darwin':
                success = self._convert_with_macos_tools(presentation_path, pdf_output_path)
            
            if success and pdf_output_path.exists():
                print(f"    ✓ Converted presentation to PDF: {presentation_path.name} → {pdf_filename}")
                return pdf_output_path
            else:
                print(f"    ⚠ Could not convert {presentation_path.name} to PDF, copying original file")
                # Copy original file to documents folder with group prefix
                original_copy = output_dir / f"{group_name}-{presentation_path.name}"
                shutil.copy2(presentation_path, original_copy)
                return original_copy
                
        except Exception as e:
            print(f"    ⚠ Error converting {presentation_path.name} to PDF: {e}")
            # Copy original file as fallback
            try:
                original_copy = output_dir / f"{group_name}-{presentation_path.name}"
                shutil.copy2(presentation_path, original_copy)
                return original_copy
            except Exception as copy_error:
                print(f"    ⚠ Error copying original file: {copy_error}")
                return presentation_path
    
    def _convert_with_libreoffice(self, input_path: Path, output_dir: Path, pdf_path: Path) -> bool:
        """Try to convert using LibreOffice headless mode."""
        try:
            # Check if LibreOffice is available (try both commands)
            soffice_cmd = None
            for cmd in ['soffice', 'libreoffice']:
                try:
                    result = subprocess.run([cmd, '--version'], 
                                          capture_output=True, text=True, timeout=10)
                    if result.returncode == 0:
                        soffice_cmd = cmd
                        break
                except FileNotFoundError:
                    continue
            
            if not soffice_cmd:
                return False
            
            # Convert using LibreOffice
            cmd = [
                soffice_cmd,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(output_dir),
                str(input_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                # LibreOffice creates PDF with same name as input file
                default_pdf = output_dir / f"{input_path.stem}.pdf"
                if default_pdf.exists():
                    # Rename to our desired filename
                    default_pdf.rename(pdf_path)
                    return True
            
            return False
            
        except (subprocess.TimeoutExpired, subprocess.SubprocessError, FileNotFoundError):
            return False
    
    def _convert_with_python_libs(self, input_path: Path, pdf_path: Path) -> bool:
        """Try to convert using Python libraries (basic conversion)."""
        try:
            # Try importing python-pptx and reportlab
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            # This is a very basic conversion - just extracts text
            prs = Presentation(str(input_path))
            
            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            y_position = 750
            
            for slide_num, slide in enumerate(prs.slides):
                if slide_num > 0:
                    c.showPage()
                    y_position = 750
                
                c.drawString(50, y_position, f"Slide {slide_num + 1}")
                y_position -= 30
                
                for shape in slide.shapes:
                    try:
                        # Try to get text from shape using different approaches
                        text = ""
                        # Use getattr to safely access attributes that may not exist
                        text_frame = getattr(shape, "text_frame", None)
                        if text_frame and hasattr(text_frame, "text"):
                            text = text_frame.text
                        elif hasattr(shape, "text"):
                            text = getattr(shape, "text", "")
                        
                        if text and text.strip():
                            lines = text.split('\n')
                            for line in lines:
                                if y_position < 50:
                                    c.showPage()
                                    y_position = 750
                                c.drawString(50, y_position, line[:80])  # Limit line length
                                y_position -= 20
                    except:
                        # Skip shapes that don't have readable text
                        continue
            
            c.save()
            return True
            
        except ImportError:
            # Libraries not available
            return False
        except Exception:
            return False
    
    def _convert_with_macos_tools(self, input_path: Path, pdf_path: Path) -> bool:
        """Try to convert using macOS-specific tools."""
        try:
            file_ext = input_path.suffix.lower()
            
            if file_ext in ['.key']:
                # Try Keynote
                applescript = f'''
                tell application "Keynote"
                    open "{input_path}"
                    export front document as PDF to file "{pdf_path}"
                    close front document
                end tell
                '''
                result = subprocess.run(['osascript', '-e', applescript], 
                                      capture_output=True, text=True, timeout=30)
                return result.returncode == 0
                
            elif file_ext in ['.ppt', '.pptx']:
                # Try Microsoft PowerPoint
                applescript = f'''
                tell application "Microsoft PowerPoint"
                    open "{input_path}"
                    save front presentation in "{pdf_path}" as save as PDF
                    close front presentation
                end tell
                '''
                result = subprocess.run(['osascript', '-e', applescript], 
                                      capture_output=True, text=True, timeout=30)
                return result.returncode == 0
            
            return False
            
        except (subprocess.TimeoutExpired, subprocess.SubprocessError):
            return False
    
    def detect_project_type(self, group_folder: Path) -> str:
        """Detect the type of project based on files and folders."""
        # Check for Unity project
        if list(group_folder.rglob("*.unity")) or list(group_folder.rglob("Assets")):
            return "Unity Project"
        
        # Check for AR/VR project indicators
        ar_indicators = ["ar", "vr", "xr", "augmented", "virtual", "mixed reality"]
        for item in group_folder.rglob("*"):
            if any(indicator in item.name.lower() for indicator in ar_indicators):
                return "AR/VR Project"
        
        # Check for web project
        if (list(group_folder.rglob("*.html")) or 
            list(group_folder.rglob("package.json")) or
            list(group_folder.rglob("index.js"))):
            return "Web Project"
        
        # Check for mobile project
        if (list(group_folder.rglob("*.xcodeproj")) or 
            list(group_folder.rglob("*.apk")) or
            list(group_folder.rglob("AndroidManifest.xml"))):
            return "Mobile Project"
        
        # Check for source code project
        source_code_exts = {'.cs', '.py', '.js', '.ts', '.cpp', '.c', '.java', '.swift'}
        if any(list(group_folder.rglob(f"*{ext}")) for ext in source_code_exts):
            return "Source Code Project"
        
        return "Document Project"
    
    def generate_organization_report(self, organized_path: Path, group_name: str, 
                                   file_types: Dict[str, List[Path]], project_type: str) -> None:
        """Generate a report of the organization process."""
        report_path = organized_path / f"{group_name}_organization_report.txt"
        
        with open(report_path, 'w') as f:
            f.write(f"Organization Report for {group_name}\n")
            f.write("="*50 + "\n\n")
            f.write(f"Project Type: {project_type}\n")
            f.write(f"Organization Date: {os.path.getmtime(organized_path)}\n\n")
            
            f.write("File Summary:\n")
            f.write("-"*20 + "\n")
            
            total_files = 0
            for category, files in file_types.items():
                if files:
                    f.write(f"{category.upper()}: {len(files)} files\n")
                    total_files += len(files)
            
            f.write(f"\nTotal Files Organized: {total_files}\n\n")
            
            f.write("Directory Structure:\n")
            f.write("-"*20 + "\n")
            f.write(f"{organized_path.name}/\n")
            f.write("├── documents/\n")
            f.write("├── videos/\n")
            f.write("└── source_code/\n")
        
        print(f"📊 Organization report saved: {report_path.name}")
    
    def clean_unnecessary_files(self, folder_path: Path) -> None:
        """Remove unnecessary files from organized folders that aren't useful for assessment."""
        if not folder_path.exists():
            return
            
        files_removed = 0
        for file_path in folder_path.rglob("*"):
            if file_path.is_file():
                # Remove Unity .meta files and other non-assessment files
                if (file_path.suffix.lower() in {'.meta', '.prefab', '.unity', '.mat', '.asset', '.anim', '.controller'} or
                    file_path.name.lower() in {'.ds_store', 'thumbs.db', 'desktop.ini'}):
                    try:
                        file_path.unlink()
                        files_removed += 1
                        print(f"  Removed unnecessary file: {file_path.name}")
                    except Exception as e:
                        print(f"  Warning: Could not remove {file_path.name}: {e}")
                        
        if files_removed > 0:
            print(f"  ✓ Cleaned {files_removed} unnecessary files from organized folder")

    def copy_source_code_selectively(self, source_folder: Path, dest_folder: Path, group_name: str) -> None:
        """Copy only assessment-useful source code files, excluding Unity metadata."""
        if not source_folder.exists():
            return
            
        copied_files = 0
        skipped_files = 0
        
        for file_path in source_folder.rglob("*"):
            if file_path.is_file() and self.is_source_code_file(file_path):
                # Calculate relative path from source folder
                relative_path = file_path.relative_to(source_folder)
                dest_file_path = dest_folder / relative_path
                
                # Create parent directories if needed
                dest_file_path.parent.mkdir(parents=True, exist_ok=True)
                
                try:
                    shutil.copy2(file_path, dest_file_path)
                    copied_files += 1
                except Exception as e:
                    print(f"  Warning: Could not copy {file_path.name}: {e}")
            elif file_path.is_file():
                skipped_files += 1
                
        print(f"  ✓ Copied {copied_files} source files, skipped {skipped_files} unnecessary files")
    
    def organize_all_groups(self, course_dir: Path) -> None:
        """Organize all group folders found in the course directory."""
        print(f"\n🗂 Now organizing groups in: {course_dir}")
        
        # Find group folders using the robust pattern matching
        all_subdirs = [d for d in course_dir.iterdir() 
                      if d.is_dir() and not d.name.startswith('.')]
        
        group_folders = [d for d in all_subdirs if is_group_folder(d.name)]
        
        print(f"Found {len(group_folders)} groups to organize: {[f.name for f in group_folders]}")
        
        # Organize each group
        for group_folder in group_folders:
            try:
                self.organize_group_folder(group_folder, course_dir)
            except Exception as e:
                print(f"Error organizing group {group_folder.name}: {e}")
                continue
        
        print(f"\n✅ Organization complete! Organized {len(group_folders)} groups.")
        print(f"Organized files are in: {course_dir / 'organized_assignments'}")

# CONFIGURATION - Edit this path to point to your zip file
ZIP_FILE_PATH = "/Users/mac/Library/CloudStorage/OneDrive-HanoiUniversityofScienceandTechnology/Study/Research/AIoT-Lab/grading-system/3702ICT_AR.zip"

def is_group_folder(folder_name: str) -> bool:
    """Check if a folder name matches typical group folder patterns."""
    import re
    
    # Common group folder patterns:
    # GC1, GC2, GC10, GC4_1, GC4_2, etc. (Group Course)
    # OL1, OL2, etc. (Online)
    # Group1, Group2, etc.
    # Team1, Team2, etc.
    patterns = [
        r'^GC\d+(_\d+)?$',   # GC1, GC2, GC10, GC4_1, GC4_2, etc.
        r'^OL\d+(_\d+)?$',   # OL1, OL2, OL1_1, etc.
        r'^Group\d+$',       # Group1, Group2, etc.
        r'^Team\d+$',        # Team1, Team2, etc.
        r'^G\d+$',           # G1, G2, etc.
        r'^T\d+$',           # T1, T2, etc.
        r'^Group[\s\-_]*\d+$',  # "Group 1", "Group-1", "Group_1", etc.
        r'^Team[\s\-_]*\d+$',   # "Team 1", "Team-1", "Team_1", etc.
    ]
    
    # Clean folder name - remove extra spaces and normalize
    clean_name = folder_name.strip()
    
    for pattern in patterns:
        if re.match(pattern, clean_name, re.IGNORECASE):
            return True
    
    # Additional check for common variations
    # Check if folder name contains "group" or "team" followed by numbers
    if re.search(r'(group|team)[\s\-_]*\d+', clean_name, re.IGNORECASE):
        return True
    
    # Check for patterns like "GC1 - Something" or "OL1_Assignment"
    if re.search(r'^(GC|OL)\d+[\s\-_]', clean_name, re.IGNORECASE):
        return True
    
    return False

def find_group_folders_directory(extracted_dir: str) -> str:
    """
    Recursively find the directory that contains group folders.
    This handles nested directory structures like 3702ICT_AR/3702ICT_AR/GC1/...
    """
    import os
    
    def find_groups_recursive(current_dir: str, max_depth: int = 4) -> str:
        """Recursively search for directory containing group folders."""
        if max_depth <= 0:
            return current_dir
        
        try:
            # Check if current directory contains group folders
            subdirs = [d for d in os.listdir(current_dir) 
                      if os.path.isdir(os.path.join(current_dir, d)) and not d.startswith('.')]
            
            # Count how many subdirectories look like group folders
            group_folder_count = sum(1 for d in subdirs if is_group_folder(d))
            
            if group_folder_count > 0:
                print(f"Found {group_folder_count} group folders in: {current_dir}")
                return current_dir
            
            # If no group folders found here, check subdirectories
            # Sort subdirectories to prioritize course-like names
            subdirs.sort(key=lambda x: (
                # Prioritize directories that contain course codes or similar names
                0 if any(keyword in x.lower() for keyword in ['ict', 'course', 'assignment']) else 1,
                x
            ))
            
            for subdir in subdirs:
                subdir_path = os.path.join(current_dir, subdir)
                try:
                    # Recursively check this subdirectory
                    deeper_result = find_groups_recursive(subdir_path, max_depth - 1)
                    
                    # Check if the deeper search found group folders
                    if deeper_result != subdir_path:  # If deeper search found something different
                        return deeper_result
                    
                    # Check the current result for group folders
                    deeper_subdirs = [d for d in os.listdir(deeper_result) 
                                    if os.path.isdir(os.path.join(deeper_result, d)) and not d.startswith('.')]
                    deeper_group_count = sum(1 for d in deeper_subdirs if is_group_folder(d))
                    
                    if deeper_group_count > 0:
                        return deeper_result
                        
                except (OSError, PermissionError):
                    continue
            
        except (OSError, PermissionError):
            pass
        
        return current_dir
    
    result = find_groups_recursive(extracted_dir)
    print(f"Using course directory: {result}")
    return result

def main():
    """Run the assignment organizer on the specified zip file."""
    # Get the current directory
    current_dir = os.getcwd()
    print(f"Working directory: {current_dir}")
    
    # Check if zip file exists
    if not os.path.exists(ZIP_FILE_PATH):
        print(f"Error: Zip file '{ZIP_FILE_PATH}' not found in current directory")
        print("Please update the ZIP_FILE_PATH variable in this script")
        return
    
    # Create organizer and process the zip file
    organizer = AssignmentOrganizer(current_dir)
    organizer.process_course_zip(ZIP_FILE_PATH)
    
    # After extraction, find the extracted folder and organize groups
    zip_name = os.path.splitext(os.path.basename(ZIP_FILE_PATH))[0]  # Get just the filename without path
    extracted_dir = os.path.join(current_dir, zip_name)
    
    if os.path.exists(extracted_dir):
        print(f"\n🗂 Now organizing groups...")
        
        # Smart detection of course directory structure
        course_dir = find_group_folders_directory(extracted_dir)
        
        # Find group folders (looking for typical group naming patterns like GC1, GC2, OL1, etc.)
        group_folders = [d for d in os.listdir(course_dir) 
                        if os.path.isdir(os.path.join(course_dir, d)) and not d.startswith('.') 
                        and is_group_folder(d)]
        
        print(f"Found {len(group_folders)} groups to organize: {group_folders}")
        
        # Extract and organize each group
        from pathlib import Path
        for group_name in group_folders:
            group_path = Path(course_dir) / group_name
            
            # First extract any zip files in the group folder
            organizer.extract_group_zips(group_path)
            
            # Now organize the extracted files - organized folders will be created in course_dir
            organizer.organize_group_folder(group_path, Path(course_dir))
    else:
        print(f"Warning: Extracted directory not found at {extracted_dir}")
    
    print("\n" + "="*50)
    print("EXTRACTION AND ORGANIZATION COMPLETE!")
    print("="*50)
    print("What was done:")
    print("  1. Extracted the main course zip file")
    print("  2. Found all group folders (GC1, GC2, OL1, etc.)")
    print("  3. For each group folder:")
    print("     - RECURSIVELY extracted ALL zip files (including nested ones)")
    print("     - Deleted the zip files after extraction")
    print("  4. Created organized folders in 'organized_assignments/' directory:")
    print("     - organized_assignments/")
    print("       ├── GroupName_organized/")
    print("       │   ├── documents/     (PDF, Excel files, converted presentations)")
    print("       │   ├── videos/        (MP4, demo videos)")
    print("       │   └── source_code/   (Scripts, code folders)")
    print("       ├── GC1_organized/")
    print("       ├── GC2_organized/")
    print("       └── ... (etc.)")
    print("\n🔄 PRESENTATION CONVERSION:")
    print("     - PowerPoint/Keynote files (.ppt/.pptx/.key) converted to PDF")
    print("     - Converted presentations placed in documents/ folder")
    print("     - No separate presentations/ folder created")
    print("\n🗜️ RECURSIVE ZIP EXTRACTION:")
    print("     - All zip files extracted recursively until no more zip files found")
    print("     - Handles nested zip files (zip files inside zip files)")
    print("     - Up to 10 extraction iterations to ensure complete extraction")
    print("\nAll files have been renamed with group prefixes!")
    print("Original extracted files are preserved in the course folder.")
    print("Organized files are in 'organized_assignments/' directory with separate folders for each group.")


if __name__ == "__main__":
    main()
